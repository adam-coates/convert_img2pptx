#!/usr/bin/env python3

import os
import argparse
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def images_to_pptx(image_dir, output_pptx):
    # Create a presentation object with widescreen dimensions
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Get a list of all image files in the directory and sort them alphabetically
    image_files = sorted([f for f in os.listdir(image_dir) if os.path.isfile(os.path.join(image_dir, f)) and f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))])

    # Add each image as a slide in the presentation
    for image_file in image_files:
        # Load the image to get its size
        img_path = os.path.join(image_dir, image_file)
        img = Image.open(img_path)
        img_width, img_height = img.size
        aspect_ratio = img_width / img_height

        # Create a new slide with dimensions matching the aspect ratio of the image
        slide_width = prs.slide_width
        slide_height = slide_width / aspect_ratio
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use slide layout for images
        left = top = 0

        # Add image to slide with adjusted dimensions to fit the slide
        slide.shapes.add_picture(img_path, left, top, width=slide_width, height=slide_height)

    # Save the presentation
    prs.save(output_pptx)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert images in a directory to a PowerPoint presentation")
    parser.add_argument("image_directory", help="Directory containing images")
    parser.add_argument("output_pptx", help="Output PowerPoint file (.pptx)")
    args = parser.parse_args()

    images_to_pptx(args.image_directory, args.output_pptx)

